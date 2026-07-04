"""첨부 파일 → 레퍼런스 텍스트 추출.

사용자가 업로드한 파일(PDF, Word, PPT, Excel, 텍스트)의 본문을 뽑아
모든 LLM에게 동일한 레퍼런스 원문으로 제공한다.
"""
import io

MAX_CHARS_PER_FILE = 30_000
_MAX_ROWS_PER_SHEET = 200


def extract_file_text(filename: str, data: bytes) -> str:
    """파일 본문 텍스트를 추출한다. 실패 시 오류 설명 문자열을 반환."""
    ext = filename.lower().rsplit(".", 1)[-1] if "." in filename else ""
    try:
        if ext == "pdf":
            text = _from_pdf(data)
        elif ext == "docx":
            text = _from_docx(data)
        elif ext == "pptx":
            text = _from_pptx(data)
        elif ext in ("xlsx", "xlsm"):
            text = _from_xlsx(data)
        elif ext in ("txt", "md", "csv"):
            text = _from_plain(data)
        else:
            return f"[지원하지 않는 파일 형식: .{ext}]"
    except Exception as e:
        return f"[추출 실패: {e}]"

    text = text.strip()
    if not text:
        return "[본문 텍스트를 찾지 못함 — 스캔 이미지 PDF일 수 있습니다]"
    if len(text) > MAX_CHARS_PER_FILE:
        text = text[:MAX_CHARS_PER_FILE] + " …(이하 생략)"
    return text


def _from_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = []
    for i, page in enumerate(reader.pages, 1):
        page_text = (page.extract_text() or "").strip()
        if page_text:
            pages.append(f"[p.{i}] {page_text}")
    return "\n".join(pages)


def _from_docx(data: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _from_pptx(data: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                texts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        texts.append(" | ".join(cells))
        if texts:
            parts.append(f"[슬라이드 {i}]\n" + "\n".join(texts))
    return "\n\n".join(parts)


def _from_xlsx(data: bytes) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        rows = []
        for r, row in enumerate(ws.iter_rows(values_only=True)):
            if r >= _MAX_ROWS_PER_SHEET:
                rows.append("…(이하 행 생략)")
                break
            cells = [str(v) for v in row if v is not None and str(v).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            parts.append(f"[시트: {ws.title}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(parts)


def _from_plain(data: bytes) -> str:
    for enc in ("utf-8", "cp949"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")
