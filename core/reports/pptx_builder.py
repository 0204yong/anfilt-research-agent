"""보고서 JSON → PowerPoint(.pptx) 변환."""
import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY = RGBColor(0x1F, 0x38, 0x64)
GREEN = RGBColor(0x2E, 0x7D, 0x5B)
GRAY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "맑은 고딕"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def _set_run(run, size, bold=False, color=NAVY):
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def _add_header(slide, title):
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.9))  # 1 = rectangle
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_run(run, 24, bold=True, color=WHITE)


def _add_bullets(slide, items, top=Inches(1.2), size=16):
    tf = _textbox(slide, Inches(0.6), top, SLIDE_W - Inches(1.2), SLIDE_H - top - Inches(0.4))
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "▪ " + str(item)
        _set_run(run, size, color=RGBColor(0x33, 0x33, 0x33))
    return tf


def _add_body_text(slide, text, top=Inches(1.2), size=14):
    tf = _textbox(slide, Inches(0.6), top, SLIDE_W - Inches(1.2), SLIDE_H - top - Inches(0.4))
    first = True
    for para in str(text).split("\n"):
        if not para.strip():
            continue
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = para.strip()
        _set_run(run, size, color=RGBColor(0x33, 0x33, 0x33))
    return tf


def _add_table_slide(prs, table_data):
    slide = _blank_slide(prs)
    _add_header(slide, table_data.get("title", "데이터"))
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])
    if not headers:
        return
    n_rows = min(len(rows), 14) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(0.6), Inches(1.2), SLIDE_W - Inches(1.2),
        Inches(0.4) * n_rows,
    )
    table = shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = str(h)
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                _set_run(run, 12, bold=True, color=WHITE)
    for r, row in enumerate(rows[:14], start=1):
        for c in range(n_cols):
            val = str(row[c]) if c < len(row) else ""
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_run(run, 11, color=RGBColor(0x33, 0x33, 0x33))


def _add_labeled_list(tf, label, items, item_size=13, first=False):
    """텍스트 프레임에 '소제목 + 불릿 목록' 블록을 추가한다."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(0 if first else 14)
    p.space_after = Pt(6)
    run = p.add_run()
    run.text = label
    _set_run(run, item_size + 2, bold=True, color=NAVY)
    for item in items:
        p = tf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = "▪ " + str(item)
        _set_run(run, item_size, color=RGBColor(0x33, 0x33, 0x33))


def _add_cover_band(slide, title, subtitle, height=Inches(1.4), title_size=26):
    band = slide.shapes.add_shape(1, 0, 0, SLIDE_W, height)
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    _set_run(run, title_size, bold=True, color=WHITE)
    if subtitle:
        p2 = tf.add_paragraph()
        run2 = p2.add_run()
        run2.text = subtitle
        _set_run(run2, 13, color=RGBColor(0xD9, 0xE2, 0xF3))


def _build_onepager(report: dict) -> bytes:
    """1장짜리 원페이저 — 제목·요약·핵심발견·제언을 한 슬라이드에 압축."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = _blank_slide(prs)
    _add_cover_band(slide, report.get("title", "조사 보고서"),
                    report.get("subtitle", ""), height=Inches(1.2), title_size=24)

    # 왼쪽: 요약 / 오른쪽: 핵심 발견 + 제언 (2단 구성)
    left_tf = _textbox(slide, Inches(0.5), Inches(1.5), Inches(6.2), Inches(5.6))
    p = left_tf.paragraphs[0]
    run = p.add_run()
    run.text = "Executive Summary"
    _set_run(run, 15, bold=True, color=GREEN)
    for para in str(report.get("executive_summary", "")).split("\n"):
        if not para.strip():
            continue
        p = left_tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = para.strip()
        _set_run(run, 13, color=RGBColor(0x33, 0x33, 0x33))

    right_tf = _textbox(slide, Inches(7.0), Inches(1.5), Inches(5.8), Inches(5.6))
    _add_labeled_list(right_tf, "핵심 발견사항", report.get("key_findings", []),
                      item_size=12, first=True)
    if report.get("recommendations"):
        _add_labeled_list(right_tf, "제언", report["recommendations"], item_size=12)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _build_compact(report: dict) -> bytes:
    """2~7장 컴팩트 구성 — [표지+요약] + 섹션/표 + [핵심발견+제언]."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1) 표지 + 요약 통합
    slide = _blank_slide(prs)
    _add_cover_band(slide, report.get("title", "조사 보고서"), report.get("subtitle", ""))
    tf = _textbox(slide, Inches(0.6), Inches(1.7), SLIDE_W - Inches(1.2), Inches(5.2))
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Executive Summary"
    _set_run(run, 16, bold=True, color=GREEN)
    for para in str(report.get("executive_summary", "")).split("\n"):
        if not para.strip():
            continue
        p = tf.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = para.strip()
        _set_run(run, 14, color=RGBColor(0x33, 0x33, 0x33))

    # 2) 본문 섹션
    for sec in report.get("sections", []):
        slide = _blank_slide(prs)
        _add_header(slide, sec.get("heading", ""))
        bullets = sec.get("bullets") or []
        if bullets:
            _add_bullets(slide, bullets, size=15)
        else:
            _add_body_text(slide, sec.get("content", ""))

    # 3) 데이터 표
    for table_data in report.get("data_tables", []):
        _add_table_slide(prs, table_data)

    # 4) 핵심 발견사항 + 제언 통합
    slide = _blank_slide(prs)
    _add_header(slide, "핵심 발견사항 · 제언")
    tf = _textbox(slide, Inches(0.6), Inches(1.2), SLIDE_W - Inches(1.2), Inches(5.9))
    _add_labeled_list(tf, "핵심 발견사항", report.get("key_findings", []), first=True)
    if report.get("recommendations"):
        _add_labeled_list(tf, "제언", report["recommendations"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_pptx(report: dict, target_pages: int = None) -> bytes:
    if target_pages is not None and target_pages <= 1:
        return _build_onepager(report)
    if target_pages is not None and target_pages <= 7:
        return _build_compact(report)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1) 표지
    slide = _blank_slide(prs)
    band = slide.shapes.add_shape(1, 0, Inches(2.6), SLIDE_W, Inches(2.2))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tf = band.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = report.get("title", "조사 보고서")
    _set_run(run, 36, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = report.get("subtitle", "")
    _set_run(run2, 18, color=RGBColor(0xD9, 0xE2, 0xF3))

    # 2) 경영진 요약
    slide = _blank_slide(prs)
    _add_header(slide, "Executive Summary")
    _add_body_text(slide, report.get("executive_summary", ""), size=16)

    # 3) 핵심 발견사항
    if report.get("key_findings"):
        slide = _blank_slide(prs)
        _add_header(slide, "핵심 발견사항")
        _add_bullets(slide, report["key_findings"])

    # 4) 본문 섹션
    for sec in report.get("sections", []):
        slide = _blank_slide(prs)
        _add_header(slide, sec.get("heading", ""))
        bullets = sec.get("bullets") or []
        if bullets:
            _add_bullets(slide, bullets, size=15)
        else:
            _add_body_text(slide, sec.get("content", ""))

    # 5) 데이터 표
    for table_data in report.get("data_tables", []):
        _add_table_slide(prs, table_data)

    # 6) 제언
    if report.get("recommendations"):
        slide = _blank_slide(prs)
        _add_header(slide, "제언 (Recommendations)")
        _add_bullets(slide, report["recommendations"])

    # 7) 출처
    if report.get("sources"):
        slide = _blank_slide(prs)
        _add_header(slide, "출처 (Sources)")
        items = [
            f"{s.get('title', '')} — {s.get('url', '')}"
            for s in report["sources"]
        ]
        _add_bullets(slide, items, size=12)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
